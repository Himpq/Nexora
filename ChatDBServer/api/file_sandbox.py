import json
import io
import os
import re
import threading
import time
import uuid
import hashlib
from typing import Any, Dict, List, Optional, Tuple

from secure import safe_filename
from datastorage import safe_read_json, safe_write_json
from text_patch import apply_text_patch, build_preview_diff, line_separator_name
import prompts


_LOCKS: Dict[str, threading.Lock] = {}
_LOCKS_GUARD = threading.Lock()


def _get_lock(username: str) -> threading.Lock:
    with _LOCKS_GUARD:
        if username not in _LOCKS:
            _LOCKS[username] = threading.Lock()
        return _LOCKS[username]


class UserFileSandbox:
    """
    用户文件沙箱：
    - 实体文件存储在 ChatDBServer/data/temp 下（随机命名）
    - 用户索引存储在 data/users/<username>/file_sandbox.json
    - 对外统一暴露 {username}/files/{alias}
    """

    ALLOWED_TEXT_EXTS = {
        ".txt", ".md", ".py", ".c", ".h", ".hpp", ".cpp", ".cc", ".cxx",
        ".js", ".ts", ".tsx", ".jsx", ".java", ".go", ".rs", ".cs", ".php",
        ".rb", ".swift", ".kt", ".kts", ".scala", ".sh", ".bash", ".zsh",
        ".bat", ".ps1", ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg",
        ".xml", ".html", ".css", ".sql", ".csv", ".log"
    }
    PARSED_BINARY_EXTS = {".docx", ".pdf", ".pptx"}
    IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
    ALLOWED_UPLOAD_EXTS = ALLOWED_TEXT_EXTS | PARSED_BINARY_EXTS | IMAGE_EXTS
    FILE_READ_MAX_LINES = 500
    FILE_READ_MAX_CHARS = 10000

    def __init__(self, username: str):
        self.username = str(username or "").strip()
        base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        self.base_dir = base_dir
        self.temp_dir = os.path.join(base_dir, "data", "temp")
        self.user_dir = os.path.join(base_dir, "data", "users", self.username)
        self.index_path = os.path.join(self.user_dir, "file_sandbox.json")
        self._lock = _get_lock(self.username)
        os.makedirs(self.temp_dir, exist_ok=True)
        os.makedirs(self.user_dir, exist_ok=True)
        self._ensure_index()

    def _ensure_index(self) -> None:
        if os.path.exists(self.index_path):
            return
        safe_write_json(self.index_path, {"files": {}}, indent=2)

    def _load_index(self) -> Dict[str, Any]:
        data = safe_read_json(self.index_path, default={})
        if not isinstance(data, dict):
            data = {}
        files = data.get("files")
        if not isinstance(files, dict):
            files = {}
        data["files"] = files
        return data

    def _save_index(self, data: Dict[str, Any]) -> None:
        safe_write_json(self.index_path, data, indent=2)

    def _sanitize_alias_directory(self, path: str) -> str:
        raw = str(path or "").replace("\\", "/").strip().strip("/")

        if not raw:
            return ""

        parts = []

        for segment in raw.split("/"):
            piece = str(segment or "").strip()

            if not piece:
                continue

            if piece in {".", ".."}:
                raise ValueError("invalid file directory")

            parts.append(safe_filename(piece, default="folder", max_len=80))

            if len(parts) > 6:
                raise ValueError("file directory cannot exceed 6 levels")

        return "/".join(parts)

    def _sanitize_alias(self, name: str) -> str:
        raw = str(name or "").replace("\\", "/").strip().strip("/")

        if not raw:
            return "untitled.txt"

        segments = raw.split("/")
        filename = segments[-1]
        directory = self._sanitize_alias_directory("/".join(segments[:-1]))
        safe_name = safe_filename(filename, default="untitled.txt", max_len=120)
        alias = f"{directory}/{safe_name}" if directory else safe_name

        if len(alias) > 260:
            raise ValueError("file alias cannot exceed 260 characters")

        return alias

    def _unique_alias(self, desired_name: str, files: Dict[str, Any]) -> str:
        alias = self._sanitize_alias(desired_name)
        if alias not in files:
            return alias
        root, ext = os.path.splitext(alias)
        i = 2
        while True:
            candidate = f"{root}-{i}{ext}"
            if candidate not in files:
                return candidate
            i += 1

    def _decode_bytes(self, b: bytes) -> Tuple[str, str]:
        encodings = ["utf-8", "utf-8-sig", "gb18030", "gbk", "latin-1"]
        for enc in encodings:
            try:
                return b.decode(enc), enc
            except Exception:
                continue
        return b.decode("utf-8", errors="replace"), "utf-8-replace"

    def _resolve_alias(self, file_ref: str) -> str:
        ref = str(file_ref or "").replace("\\", "/").strip().strip("/")
        prefix = f"{self.username}/files/"
        alt_prefix = f"files/{self.username}/"

        if ref.startswith(prefix):
            ref = ref[len(prefix):]

        if ref.startswith(alt_prefix):
            ref = ref[len(alt_prefix):]

        if not ref:
            raise ValueError("file_ref is empty")

        return self._sanitize_alias(ref)

    def _extract_text_from_docx(self, b: bytes) -> str:
        try:
            from docx import Document  # type: ignore
        except Exception:
            raise RuntimeError("缺少依赖 python-docx，请安装: pip install python-docx")

        try:
            doc = Document(io.BytesIO(b))
            chunks: List[str] = []

            for p in doc.paragraphs:
                text = (p.text or "").strip()
                if text:
                    chunks.append(text)

            for t_idx, table in enumerate(doc.tables, start=1):
                chunks.append(f"\n[Table {t_idx}]")
                for row in table.rows:
                    cells = []
                    for cell in row.cells:
                        cell_text = (cell.text or "").replace("\n", " ").strip()
                        cells.append(cell_text)
                    line = " | ".join(cells).strip()
                    if line:
                        chunks.append(line)

            result = "\n".join(chunks).strip()
            if not result:
                raise RuntimeError("DOCX 文件未提取到可读文本（可能仅含图片/公式）")
            return result
        except RuntimeError:
            raise
        except Exception as e:
            raise RuntimeError(f"DOCX 解析失败: {e}")

    def _extract_text_from_pdf(self, b: bytes) -> str:
        pdf_reader_cls = None
        try:
            from pypdf import PdfReader  # type: ignore
            pdf_reader_cls = PdfReader
        except Exception:
            try:
                from PyPDF2 import PdfReader  # type: ignore
                pdf_reader_cls = PdfReader
            except Exception:
                raise RuntimeError("缺少 PDF 解析依赖，请安装: pip install pypdf")

        try:
            reader = pdf_reader_cls(io.BytesIO(b))
            chunks: List[str] = []
            for i, page in enumerate(reader.pages, start=1):
                page_text = page.extract_text() or ""
                page_text = page_text.strip()
                if page_text:
                    chunks.append(f"[Page {i}]")
                    chunks.append(page_text)

            result = "\n\n".join(chunks).strip()
            if not result:
                raise RuntimeError("PDF 未提取到可读文本（可能是扫描件图片 PDF）")
            return result
        except RuntimeError:
            raise
        except Exception as e:
            raise RuntimeError(f"PDF 解析失败: {e}")

    def _extract_text_from_pptx(self, b: bytes) -> str:
        try:
            from pptx import Presentation  # type: ignore
        except Exception:
            raise RuntimeError("缺少依赖 python-pptx，请安装: pip install python-pptx")

        try:
            prs = Presentation(io.BytesIO(b))
            chunks: List[str] = []
            for slide_idx, slide in enumerate(prs.slides, start=1):
                slide_lines: List[str] = []
                table_count = 0

                for shape in slide.shapes:
                    if getattr(shape, "has_table", False):
                        table_count += 1
                        slide_lines.append(f"[Table {table_count}]")
                        try:
                            table = shape.table
                            for row in table.rows:
                                cells = []
                                for cell in row.cells:
                                    t = (cell.text or "").replace("\n", " ").strip()
                                    cells.append(t)
                                line = " | ".join(cells).strip()
                                if line:
                                    slide_lines.append(line)
                        except Exception:
                            pass
                        continue

                    text = ""
                    if hasattr(shape, "text_frame") and getattr(shape, "text_frame", None):
                        text = shape.text_frame.text or ""
                    elif hasattr(shape, "text"):
                        text = getattr(shape, "text", "") or ""
                    text = text.replace("\x0b", "\n").strip()
                    if text:
                        slide_lines.append(text)

                # 提取备注区文本
                try:
                    notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                    if notes:
                        slide_lines.append("[Notes]")
                        slide_lines.append(notes)
                except Exception:
                    pass

                if slide_lines:
                    chunks.append(f"[Slide {slide_idx}]")
                    chunks.extend(slide_lines)

            result = "\n".join(chunks).strip()
            if not result:
                raise RuntimeError("PPTX 未提取到可读文本（可能仅含图片）")
            return result
        except RuntimeError:
            raise
        except Exception as e:
            raise RuntimeError(f"PPTX 解析失败: {e}")

    def _extract_upload_text(self, file_bytes: bytes, ext: str) -> Tuple[str, str, str]:
        ext = str(ext or "").lower()
        if ext in self.ALLOWED_TEXT_EXTS:
            text, enc = self._decode_bytes(bytes(file_bytes))
            return text, enc, "text"
        if ext == ".docx":
            text = self._extract_text_from_docx(bytes(file_bytes))
            return text, "docx-extract", "docx"
        if ext == ".pdf":
            text = self._extract_text_from_pdf(bytes(file_bytes))
            return text, "pdf-extract", "pdf"
        if ext == ".pptx":
            text = self._extract_text_from_pptx(bytes(file_bytes))
            return text, "pptx-extract", "pptx"
        raise ValueError(f"unsupported file extension: {ext}")

    def add_upload(
        self,
        file_bytes: bytes,
        original_name: str,
        update_file_name: Optional[str] = None,
        target_path: Optional[str] = None,
    ) -> Dict[str, Any]:
        if not isinstance(file_bytes, (bytes, bytearray)):
            raise ValueError("file_bytes must be bytes")
        original_name = os.path.basename(str(original_name or "").strip())
        if not original_name:
            original_name = "untitled.txt"
        ext = os.path.splitext(original_name)[1].lower()
        if ext not in self.ALLOWED_UPLOAD_EXTS:
            raise ValueError(f"unsupported file extension: {ext}")

        is_image = ext in self.IMAGE_EXTS
        detected_encoding = "binary" if is_image else ""
        parser_mode = "image" if is_image else ""
        text = "" if is_image else None

        if not is_image:
            text, detected_encoding, parser_mode = self._extract_upload_text(bytes(file_bytes), ext)

        # docx/pdf/pptx 解析后统一存为可编辑文本文件
        storage_ext = ext if ext in (self.ALLOWED_TEXT_EXTS | self.IMAGE_EXTS) else ".md"
        random_name = f"{uuid.uuid4().hex}{storage_ext}"
        stored_path = os.path.join(self.temp_dir, random_name)

        with self._lock:
            index = self._load_index()
            files = index.get("files", {})
            target_dir = self._sanitize_alias_directory(str(target_path or ""))
            desired_name = update_file_name or original_name
            desired_alias = f"{target_dir}/{desired_name}" if target_dir else desired_name
            alias = self._unique_alias(desired_alias, files)

            if is_image:
                with open(stored_path, "wb") as f:
                    f.write(bytes(file_bytes))
            else:
                with open(stored_path, "w", encoding="utf-8") as f:
                    f.write(str(text or ""))

            now = int(time.time())
            entry = {
                "alias": alias,
                "original_name": original_name,
                "stored_name": random_name,
                "stored_path": os.path.relpath(stored_path, self.base_dir).replace("\\", "/"),
                "sandbox_path": f"{self.username}/files/{alias}",
                "encoding": detected_encoding,
                "source_ext": ext,
                "parser_mode": parser_mode,
                "created_at": now,
                "updated_at": now,
                "size": os.path.getsize(stored_path),
            }
            files[alias] = entry
            index["files"] = files
            self._save_index(index)
            return entry

    def create_file(self, file_ref: str, content: Optional[str] = "", overwrite: bool = False) -> Dict[str, Any]:
        alias = self._resolve_alias(file_ref)
        alias = self._sanitize_alias(alias)
        if not alias:
            raise ValueError("invalid file_path")

        ext = os.path.splitext(alias)[1].lower()
        if not ext:
            alias = f"{alias}.txt"
            ext = ".txt"
        if ext not in self.ALLOWED_TEXT_EXTS:
            raise ValueError(f"cloud_file_create 仅支持文本后缀，当前不支持: {ext}")

        text = "" if content is None else str(content)
        now = int(time.time())

        with self._lock:
            index = self._load_index()
            files = index.get("files", {})
            existing = files.get(alias)

            if existing:
                if not bool(overwrite):
                    raise FileExistsError(f"file already exists: {self.username}/files/{alias}")

                try:
                    abs_path = self._get_abs_path(existing)
                except Exception:
                    random_name = f"{uuid.uuid4().hex}{ext}"
                    abs_path = os.path.join(self.temp_dir, random_name)
                    existing["stored_name"] = random_name
                    existing["stored_path"] = os.path.relpath(abs_path, self.base_dir).replace("\\", "/")

                with open(abs_path, "w", encoding="utf-8") as f:
                    f.write(text)

                existing["encoding"] = "utf-8"
                existing["source_ext"] = ext
                existing["parser_mode"] = "created"
                existing["updated_at"] = now
                existing["size"] = os.path.getsize(abs_path)
                files[alias] = existing
                index["files"] = files
                self._save_index(index)
                return {
                    "success": True,
                    "created": False,
                    "overwritten": True,
                    "file": {
                        "alias": existing.get("alias"),
                        "sandbox_path": existing.get("sandbox_path"),
                        "original_name": existing.get("original_name"),
                        "size": existing.get("size"),
                        "updated_at": existing.get("updated_at"),
                    },
                }

            random_name = f"{uuid.uuid4().hex}{ext}"
            stored_path = os.path.join(self.temp_dir, random_name)
            with open(stored_path, "w", encoding="utf-8") as f:
                f.write(text)

            entry = {
                "alias": alias,
                "original_name": alias,
                "stored_name": random_name,
                "stored_path": os.path.relpath(stored_path, self.base_dir).replace("\\", "/"),
                "sandbox_path": f"{self.username}/files/{alias}",
                "encoding": "utf-8",
                "source_ext": ext,
                "parser_mode": "created",
                "created_at": now,
                "updated_at": now,
                "size": os.path.getsize(stored_path),
            }
            files[alias] = entry
            index["files"] = files
            self._save_index(index)
            return {
                "success": True,
                "created": True,
                "overwritten": False,
                "file": {
                    "alias": entry.get("alias"),
                    "sandbox_path": entry.get("sandbox_path"),
                    "original_name": entry.get("original_name"),
                    "size": entry.get("size"),
                    "created_at": entry.get("created_at"),
                },
            }

    def write_docx_file(
        self,
        file_ref: str,
        docx_bytes: bytes,
        markdown: str,
        title: str = "",
        overwrite: bool = False,
        render_stats: Optional[Dict[str, Any]] = None,
    ) -> Dict[str, Any]:
        """保存后端生成的 Word 二进制文件，并登记到用户云文件沙箱。"""
        alias = self._resolve_alias(file_ref)
        alias = self._sanitize_alias(alias)

        if not alias:
            raise ValueError("invalid file_path")

        root, ext = os.path.splitext(alias)

        if not ext:
            alias = f"{alias}.docx"
            ext = ".docx"

        if ext.lower() != ".docx":
            raise ValueError("cloud_doc_write 只支持 .docx 文件路径")

        if not isinstance(docx_bytes, (bytes, bytearray)) or not docx_bytes:
            raise ValueError("docx_bytes is empty")

        stats = render_stats if isinstance(render_stats, dict) else {}
        now = int(time.time())
        random_name = f"{uuid.uuid4().hex}.docx"
        stored_path = os.path.join(self.temp_dir, random_name)

        with self._lock:
            index = self._load_index()
            files = index.get("files", {})
            existing = files.get(alias)
            old_abs_path = self._get_abs_path(existing) if existing else ""

            if existing and not bool(overwrite):
                raise FileExistsError(f"file already exists: {self.username}/files/{alias}")

            with open(stored_path, "wb") as f:
                f.write(bytes(docx_bytes))

            if old_abs_path and os.path.normcase(old_abs_path) != os.path.normcase(stored_path):
                os.remove(old_abs_path)

            entry = existing if existing else {}
            created_at = int(entry.get("created_at") or now)
            entry.update({
                "alias": alias,
                "original_name": alias,
                "stored_name": random_name,
                "stored_path": os.path.relpath(stored_path, self.base_dir).replace("\\", "/"),
                "sandbox_path": f"{self.username}/files/{alias}",
                "encoding": "binary",
                "source_ext": ".docx",
                "parser_mode": "docx-generated",
                "title": str(title or "").strip(),
                "markdown_chars": len(str(markdown or "")),
                "block_count": int(stats.get("block_count") or 0),
                "created_at": created_at,
                "updated_at": now,
                "size": os.path.getsize(stored_path),
            })
            files[alias] = entry
            index["files"] = files
            self._save_index(index)

            return {
                "success": True,
                "created": not bool(existing),
                "overwritten": bool(existing),
                "mode": "docx_generate",
                "markdown_chars": len(str(markdown or "")),
                "block_count": int(stats.get("block_count") or 0),
                "file": {
                    "alias": entry.get("alias"),
                    "sandbox_path": entry.get("sandbox_path"),
                    "original_name": entry.get("original_name"),
                    "size": entry.get("size"),
                    "created_at": entry.get("created_at"),
                    "updated_at": entry.get("updated_at"),
                },
            }

    def list_files(
        self,
        query: Optional[str] = None,
        regex: bool = False,
        offset: int = 0,
        limit: int = 200,
    ) -> Dict[str, Any]:
        with self._lock:
            files = self._load_index().get("files", {})
            items = list(files.values())

        q = str(query or "").strip()
        if q:
            filtered = []
            if regex:
                try:
                    pat = re.compile(q, re.IGNORECASE)
                except Exception:
                    raise ValueError(f"invalid regex: {q}")
                for x in items:
                    target = f"{x.get('alias', '')} {x.get('original_name', '')} {x.get('sandbox_path', '')}"
                    if pat.search(target):
                        filtered.append(x)
            else:
                ql = q.lower()
                for x in items:
                    target = f"{x.get('alias', '')} {x.get('original_name', '')} {x.get('sandbox_path', '')}".lower()
                    if ql in target:
                        filtered.append(x)
            items = filtered

        items.sort(key=lambda x: int(x.get("updated_at", 0) or 0), reverse=True)
        total = len(items)
        safe_offset = max(0, min(int(offset or 0), total))
        safe_limit = max(1, min(int(limit or 200), 1000))

        return {
            "total": total,
            "offset": safe_offset,
            "limit": safe_limit,
            "files": items[safe_offset:safe_offset + safe_limit],
        }

    def _get_entry(self, file_ref: str) -> Dict[str, Any]:
        alias = self._resolve_alias(file_ref)
        with self._lock:
            index = self._load_index()
            files = index.get("files", {})
            entry = files.get(alias)
            if not entry:
                raise FileNotFoundError(f"file not found: {file_ref}")
            return entry

    def _get_abs_path(self, entry: Dict[str, Any]) -> str:
        rel = str(entry.get("stored_path") or "").replace("\\", "/")
        if not rel:
            raise FileNotFoundError("stored_path missing")
        abs_path = os.path.normpath(os.path.join(self.base_dir, rel))
        if not os.path.isfile(abs_path):
            raise FileNotFoundError(f"stored file missing: {abs_path}")
        return abs_path

    def _is_generated_docx_entry(self, entry: Dict[str, Any]) -> bool:
        alias_ext = os.path.splitext(str(entry.get("alias") or ""))[1].lower()
        return alias_ext == ".docx" and str(entry.get("parser_mode") or "") == "docx-generated"

    def _is_image_entry(self, entry: Dict[str, Any]) -> bool:
        source_ext = str(entry.get("source_ext") or "").lower()
        alias_ext = os.path.splitext(str(entry.get("alias") or ""))[1].lower()
        return source_ext in self.IMAGE_EXTS or alias_ext in self.IMAGE_EXTS or str(entry.get("parser_mode") or "") == "image"

    def _reject_docx_text_mutation(self, entry: Dict[str, Any], tool_name: str) -> None:
        alias_ext = os.path.splitext(str(entry.get("alias") or ""))[1].lower()

        if alias_ext == ".docx":
            raise ValueError(f"{tool_name} 不支持写入 .docx，请使用 cloud_doc_write 重新生成 Word 文件。")

    def read_file(
        self,
        file_ref: str,
        from_line: Optional[int] = None,
        to_line: Optional[int] = None,
        from_pos: Optional[int] = None,
        to_pos: Optional[int] = None,
    ) -> Dict[str, Any]:
        entry = self._get_entry(file_ref)
        abs_path = self._get_abs_path(entry)

        if self._is_image_entry(entry):
            raise ValueError("图片文件不支持文本读取，请使用预览或下载。")

        if self._is_generated_docx_entry(entry):
            with open(abs_path, "rb") as f:
                content = self._extract_text_from_docx(f.read())
        else:
            with open(abs_path, "r", encoding="utf-8") as f:
                content = f.read()

        def _line_col_at(text: str, pos: int) -> Tuple[int, int]:
            p = max(0, min(int(pos), len(text)))
            line = 1
            col = 1
            i = 0
            while i < p:
                if text[i] == "\n":
                    line += 1
                    col = 1
                else:
                    col += 1
                i += 1
            return line, col

        total_chars = len(content)
        result_content = content
        mode = "full"
        slice_meta: Dict[str, Any] = {}
        selected_start_pos = 0

        if from_line is not None or to_line is not None:
            lines = content.splitlines(keepends=True)
            start = 1 if from_line is None else int(from_line)
            end = len(lines) if to_line is None else int(to_line)
            start = max(1, start)
            end = max(1, end)
            if end < start:
                start, end = end, start
            start_i = start - 1
            end_i = min(end, len(lines))
            result_content = "".join(lines[start_i:end_i])
            selected_start_pos = len("".join(lines[:start_i]))
            mode = "line_range"
            slice_meta = {"from_line": start, "to_line": end}
        elif from_pos is not None or to_pos is not None:
            start = 0 if from_pos is None else int(from_pos)
            end = len(content) if to_pos is None else int(to_pos)
            start = max(0, min(start, len(content)))
            end = max(0, min(end, len(content)))
            if end < start:
                start, end = end, start
            result_content = content[start:end]
            selected_start_pos = start
            mode = "char_range"
            slice_meta = {"offset": start, "length": end - start, "end_offset": end}

        # Enforce a hard read cap for model-facing outputs:
        # at most 500 lines and 10,000 chars per call.
        limit_lines = int(self.FILE_READ_MAX_LINES)
        limit_chars = int(self.FILE_READ_MAX_CHARS)
        truncated = False
        truncated_line = 0
        truncated_col = 0
        truncate_notice = ""

        effective_cut = len(result_content)
        lines_ke = result_content.splitlines(keepends=True)
        if len(lines_ke) > limit_lines:
            effective_cut = min(effective_cut, len("".join(lines_ke[:limit_lines])))
        if len(result_content) > limit_chars:
            effective_cut = min(effective_cut, limit_chars)

        if effective_cut < len(result_content):
            truncated = True
            result_content = result_content[:effective_cut]
            absolute_cut_pos = selected_start_pos + effective_cut
            truncated_line, truncated_col = _line_col_at(content, absolute_cut_pos)
            truncate_notice = prompts.build_cloud_file_read_truncate_notice(
                limit_lines,
                limit_chars,
                truncated_line,
                truncated_col,
            )
            # Put notice in content body so model sees it directly.
            result_content += truncate_notice

        return {
            "success": True,
            "mode": mode,
            "file": {
                "alias": entry.get("alias"),
                "sandbox_path": entry.get("sandbox_path"),
                "original_name": entry.get("original_name"),
                "size": entry.get("size", 0),
                "total_chars": total_chars,
            },
            "slice": slice_meta,
            "limits": {
                "max_lines": limit_lines,
                "max_chars": limit_chars,
            },
            "truncated": truncated,
            "truncate_at": {
                "line": truncated_line,
                "column": truncated_col,
            } if truncated else None,
            "truncate_notice": truncate_notice if truncated else "",
            "content": result_content,
        }

    def write_file(
        self,
        file_ref: str,
        content: Optional[str] = None,
        from_line: Optional[int] = None,
        to_line: Optional[int] = None,
        replacement: Optional[str] = None,
        old_text: Optional[str] = None,
        new_text: Optional[str] = None,
        regex: bool = False,
        max_replace: Optional[int] = None,
    ) -> Dict[str, Any]:
        entry = self._get_entry(file_ref)
        self._reject_docx_text_mutation(entry, "cloud_file_write")
        alias = str(entry.get("alias"))
        with self._lock:
            index = self._load_index()
            files = index.get("files", {})
            entry = files.get(alias)
            if not entry:
                raise FileNotFoundError(f"file not found: {file_ref}")
            self._reject_docx_text_mutation(entry, "cloud_file_write")
            abs_path = self._get_abs_path(entry)
            with open(abs_path, "r", encoding="utf-8") as f:
                current = f.read()

            mode = ""
            replaced_count = 0
            if content is not None:
                updated = str(content)
                mode = "overwrite"
            elif from_line is not None or to_line is not None:
                lines = current.splitlines(keepends=True)
                start = 1 if from_line is None else int(from_line)
                end = len(lines) if to_line is None else int(to_line)
                start = max(1, start)
                end = max(1, end)
                if end < start:
                    start, end = end, start
                start_i = start - 1
                end_i = min(end, len(lines))
                rep = "" if replacement is None else str(replacement)
                rep_lines = rep.splitlines(keepends=True)
                lines[start_i:end_i] = rep_lines
                updated = "".join(lines)
                mode = "line_replace"
                replaced_count = max(0, end_i - start_i)
            elif old_text is not None:
                old = str(old_text)
                new = "" if new_text is None else str(new_text)
                if regex:
                    flags = 0
                    count = 0 if max_replace is None else max(0, int(max_replace))
                    pat = re.compile(old, flags)
                    updated, replaced_count = pat.subn(new, current, count=count)
                else:
                    count = -1 if max_replace is None else max(0, int(max_replace))
                    if count == 0:
                        updated = current
                        replaced_count = 0
                    else:
                        updated = current.replace(old, new, count)
                        if old:
                            replaced_count = current.count(old) if count < 0 else min(current.count(old), count)
                    mode = "text_replace"
            else:
                raise ValueError("no valid write operation provided")

            with open(abs_path, "w", encoding="utf-8") as f:
                f.write(updated)

            now = int(time.time())
            entry["updated_at"] = now
            entry["size"] = os.path.getsize(abs_path)
            files[alias] = entry
            index["files"] = files
            self._save_index(index)

            return {
                "success": True,
                "mode": mode,
                "replaced_count": replaced_count,
                "file": {
                    "alias": entry.get("alias"),
                    "sandbox_path": entry.get("sandbox_path"),
                    "size": entry.get("size"),
                    "updated_at": entry.get("updated_at"),
                },
            }

    def patch_file(
        self,
        file_ref: str,
        patch: Optional[str] = None,
        edits: Optional[List[Dict[str, Any]]] = None,
        dry_run: bool = False,
        expected_sha256: Optional[str] = None,
        tool_name: str = "cloud_file_edit",
    ) -> Dict[str, Any]:
        """对云端文本文件应用统一 diff 或结构化 edits。"""
        safe_tool_name = str(tool_name or "cloud_file_edit").strip() or "cloud_file_edit"
        entry = self._get_entry(file_ref)
        self._reject_docx_text_mutation(entry, safe_tool_name)
        alias = str(entry.get("alias"))

        with self._lock:
            index = self._load_index()
            files = index.get("files", {})
            entry = files.get(alias)

            if not entry:
                raise FileNotFoundError(f"file not found: {file_ref}")

            self._reject_docx_text_mutation(entry, safe_tool_name)
            abs_path = self._get_abs_path(entry)

            try:
                with open(abs_path, "rb") as f:
                    old_raw = f.read()

                current = old_raw.decode("utf-8")
            except UnicodeDecodeError as exc:
                return {
                    "success": False,
                    "message": f"文件无法按 utf-8 解码: {exc}",
                    "path": entry.get("sandbox_path") or file_ref,
                }

            old_sha256 = hashlib.sha256(old_raw).hexdigest()
            expected = str(expected_sha256 or "").strip().lower()

            if expected and expected != old_sha256:
                return {
                    "success": False,
                    "message": "文件内容 SHA256 与 expected_sha256 不一致，已拒绝修改。",
                    "path": entry.get("sandbox_path") or file_ref,
                    "old_sha256": old_sha256,
                    "expected_sha256": expected,
                }

            updated, stats, apply_error = apply_text_patch(
                current,
                patch_text=str(patch or ""),
                edits=edits,
            )

            if apply_error:
                return {
                    "success": False,
                    "message": apply_error,
                    "path": entry.get("sandbox_path") or file_ref,
                    "old_sha256": old_sha256,
                }

            new_raw = updated.encode("utf-8")
            new_sha256 = hashlib.sha256(new_raw).hexdigest()
            diff_text = build_preview_diff(alias, current, updated)
            changed = updated != current

            if changed and not bool(dry_run):
                with open(abs_path, "wb") as f:
                    f.write(new_raw)

                now = int(time.time())
                entry["updated_at"] = now
                entry["size"] = os.path.getsize(abs_path)
                files[alias] = entry
                index["files"] = files
                self._save_index(index)

            return {
                "success": True,
                "changed": changed,
                "dry_run": bool(dry_run),
                "path": entry.get("sandbox_path") or file_ref,
                "mode": stats.get("mode", ""),
                "old_sha256": old_sha256,
                "new_sha256": new_sha256,
                "diff": diff_text,
                "line_separator": line_separator_name(updated),
                **stats,
                "file": {
                    "alias": entry.get("alias"),
                    "sandbox_path": entry.get("sandbox_path"),
                    "size": len(new_raw) if bool(dry_run) else entry.get("size"),
                    "updated_at": entry.get("updated_at"),
                },
            }

    def find_in_file(
        self,
        file_ref: str,
        keyword: str,
        regex: bool = False,
        case_sensitive: bool = True,
        max_results: int = 200,
    ) -> Dict[str, Any]:
        if not str(keyword or "").strip():
            raise ValueError("keyword is required")
        entry = self._get_entry(file_ref)
        abs_path = self._get_abs_path(entry)

        if self._is_generated_docx_entry(entry):
            with open(abs_path, "rb") as f:
                text = self._extract_text_from_docx(f.read())
        else:
            with open(abs_path, "r", encoding="utf-8") as f:
                text = f.read()

        lines = text.splitlines()
        results = []
        limit = max(1, min(int(max_results or 200), 5000))
        flags = 0 if case_sensitive else re.IGNORECASE

        if regex:
            pattern = re.compile(str(keyword), flags)
            for line_no, line in enumerate(lines, start=1):
                for m in pattern.finditer(line):
                    results.append({
                        "line": line_no,
                        "column": m.start() + 1,
                        "end_column": m.end(),
                        "match": m.group(0),
                        "text": line,
                    })
                    if len(results) >= limit:
                        break
                if len(results) >= limit:
                    break
        else:
            key = str(keyword)
            key_cmp = key if case_sensitive else key.lower()
            for line_no, line in enumerate(lines, start=1):
                source = line if case_sensitive else line.lower()
                start = 0
                while True:
                    pos = source.find(key_cmp, start)
                    if pos < 0:
                        break
                    results.append({
                        "line": line_no,
                        "column": pos + 1,
                        "end_column": pos + len(key_cmp),
                        "match": line[pos:pos + len(key_cmp)],
                        "text": line,
                    })
                    if len(results) >= limit:
                        break
                    start = pos + max(1, len(key_cmp))
                if len(results) >= limit:
                    break

        return {
            "success": True,
            "file": {
                "alias": entry.get("alias"),
                "sandbox_path": entry.get("sandbox_path"),
                "original_name": entry.get("original_name"),
            },
            "keyword": keyword,
            "regex": bool(regex),
            "case_sensitive": bool(case_sensitive),
            "matched": len(results),
            "results": results,
        }

    def remove_file(self, file_ref: str) -> Dict[str, Any]:
        alias = self._resolve_alias(file_ref)
        with self._lock:
            index = self._load_index()
            files = index.get("files", {})
            entry = files.get(alias)
            if not entry:
                return {"success": False, "message": f"file not found: {file_ref}"}
            abs_path = ""
            try:
                abs_path = self._get_abs_path(entry)
            except Exception:
                # 允许仅清理索引，兼容实体文件被外部删除的场景
                abs_path = ""
            try:
                if os.path.exists(abs_path):
                    os.remove(abs_path)
            except Exception:
                pass
            del files[alias]
            index["files"] = files
            self._save_index(index)
            return {
                "success": True,
                "removed": {
                    "alias": entry.get("alias"),
                    "sandbox_path": entry.get("sandbox_path"),
                    "original_name": entry.get("original_name"),
                },
            }
